import json, os
from pptx import Presentation

INPUT = '/root/input.pptx'
OUTPUT = '/root/result.json'

prs = Presentation(INPUT)
slides_out = []
for i, slide in enumerate(prs.slides, start=1):
    parts = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = ''.join(run.text for run in para.runs) if para.runs else para.text
            if text:
                parts.append(text)
    slides_out.append({'index': i, 'text': '\n'.join(parts)})

result = {'slide_count': len(slides_out), 'slides': slides_out}
with open(OUTPUT, 'w', encoding='utf-8') as f:
    json.dump(result, f, ensure_ascii=False, indent=2)
